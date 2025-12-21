import logging
import io
import csv
import tempfile
import os
import asyncio
import pandas as pd
from pypdf import PdfReader

# Setup logging
try:
    from backend.app.module.Functions_module import setup_logger
    logger = setup_logger()
except ImportError:
    try:
        from module.Functions_module import setup_logger
        logger = setup_logger()
    except ImportError:
        logging.basicConfig(level=logging.INFO, format='%(asctime)s %(levelname)s %(name)s %(message)s')
        logger = logging.getLogger(__name__)

class ParserService:
    """Service responsible for extracting text and links from uploaded files.

    This encapsulates the previous free-function `parse_content` and provides a
    single, testable interface for file parsing.
    """

    def __init__(self):
        self.logger = logger

    def parse_content(self, filename: str, content: bytes) -> str:

        lower = filename.lower()

        try:
            # Plain text and markdown
            if lower.endswith('.txt') or lower.endswith('.md'):
                return content.decode('utf-8', errors='replace')

            # Enhanced CSV parsing: try pandas if available for better formats
            elif lower.endswith('.csv'):
                s = content.decode('utf-8', errors='replace')
                try:

                    # Try to detect a header-oriented CSV with semicolon separators
                    lines = s.splitlines()
                    csv_start = 0
                    gruppe_navn = ''
                    found_header = False

                    for i, line in enumerate(lines):
                        line = line.strip()
                        if line.startswith('Sheet:'):
                            gruppe_navn = line.replace('Sheet:', '').strip()
                        elif ';' in line or ',' in line:
                            csv_start = i
                            found_header = True
                            break

                    csv_lines = '\n'.join(lines[csv_start:]) if found_header else s

                    with tempfile.NamedTemporaryFile(mode='w', suffix='.csv', delete=False, encoding='utf-8') as tmpf:
                        tmpf.write(csv_lines)
                        tmp_path = tmpf.name

                    try:
                        # Prefer semicolon when present
                        sep = ';' if ';' in csv_lines else ','
                        df = pd.read_csv(tmp_path, sep=sep, encoding='utf-8')
                        content_out = f"**Data from {filename}**\n\n"

                        for _, row in df.iterrows():
                            row_data = []
                            for col in df.columns:
                                val = row[col]
                                if pd.notna(val):
                                    row_data.append(f"{col}: {val}")
                            content_out += ' | '.join(row_data) + '\n'

                        try:
                            import os as _os
                            _os.unlink(tmp_path)
                        except Exception:
                            pass

                        return content_out
                    except Exception:
                        # Fallback to simple csv reader
                        rows = list(csv.reader(io.StringIO(s)))
                        lines = [', '.join(row) for row in rows]
                        return '\n'.join(lines)
                except Exception:
                    # pandas or tempfile not available - simple fallback
                    rows = list(csv.reader(io.StringIO(s)))
                    lines = [', '.join(row) for row in rows]
                    return '\n'.join(lines)

            # PDF parsing with link extraction
            if lower.endswith('.pdf'):
                try:
                    text_parts = []
                    links = []

                    try:
                        # Prefer pypdf for link extraction and robust text
                        
                        reader = PdfReader(io.BytesIO(content))
                        for i, page in enumerate(reader.pages, start=1):
                            try:
                                page_text = page.extract_text() or ''
                                if page_text.strip():
                                    text_parts.append(page_text)
                            except Exception:
                                pass

                            # Extract annotations/links if present
                            try:
                                annots_obj = page.get('/Annots')  # type: ignore[reportGeneralTypeIssues]
                                if annots_obj:
                                    # annots_obj may be an array-like PdfObject; iterate defensively
                                    try:
                                        for a in annots_obj:
                                            try:
                                                obj = a.get_object()
                                                if obj.get('/Subtype') == '/Link' and '/A' in obj:
                                                    action = obj['/A']
                                                    if action.get('/S') == '/URI' and '/URI' in action:
                                                        uri = action['/URI']
                                                        links.append(f"Link (side {i}): {uri}")
                                            except Exception:
                                                continue
                                    except TypeError:
                                        # Not iterable - try single object handling
                                        try:
                                            a = annots_obj
                                            obj = a.get_object()
                                            if obj.get('/Subtype') == '/Link' and '/A' in obj:
                                                action = obj['/A']
                                                if action.get('/S') == '/URI' and '/URI' in action:
                                                    uri = action['/URI']
                                                    links.append(f"Link (side {i}): {uri}")
                                        except Exception:
                                            pass
                            except Exception:
                                # Some pdf versions have different structures; ignore
                                pass

                    except Exception:
                        # Fallback to pdfplumber for text only (no annotations extraction)
                        try:
                            import pdfplumber
                            with pdfplumber.open(io.BytesIO(content)) as pdf:
                                for page in pdf.pages:
                                    txt = page.extract_text() or ''
                                    if txt.strip():
                                        text_parts.append(txt)
                                    else:
                                        # Try OCR fallback if pytesseract is available
                                        try:
                                            from PIL import Image
                                            import pytesseract
                                            img = page.to_image(resolution=150).original
                                            ocr_txt = pytesseract.image_to_string(img)
                                            text_parts.append(ocr_txt)
                                        except Exception:
                                            pass
                        except Exception:
                            return f"[Simulated PDF parsing for {filename} - pdf parsing libraries missing or failed]"

                    combined = '\n'.join([p for p in text_parts if p and p.strip()])

                    if links:
                        combined += '\n\n**Links from PDF document:**\n'
                        for l in links:
                            combined += f"- {l}\n"

                    if combined.strip() == '':
                        return f"[Simulated PDF parsing for {filename} - no text extracted]"

                    return combined
                except Exception as e:
                    self.logger.exception('PDF parsing error: %s', e)
                    return f"[Simulated PDF parsing for {filename} - error]"

            # DOCX parsing with link extraction
            elif lower.endswith('.docx'):
                try:
                    from docx import Document

                    doc = Document(io.BytesIO(content))
                    parts = []

                    # Extract paragraphs
                    for para in doc.paragraphs:
                        if para.text and para.text.strip():
                            parts.append(para.text.strip())

                    # Extract hyperlinks using rels
                    try:
                        all_hyperlinks = {}
                        for rel in doc.part.rels.values():
                            if 'hyperlink' in rel.reltype:
                                all_hyperlinks[rel.rId] = rel.target_ref

                        all_document_links = []
                        for paragraph in doc.paragraphs:
                            try:
                                for hyperlink in paragraph._element.findall(
                                    './/{http://schemas.openxmlformats.org/wordprocessingml/2006/main}hyperlink'
                                ):
                                    rel_id = hyperlink.get(
                                        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
                                    )
                                    if rel_id in all_hyperlinks:
                                        url = all_hyperlinks[rel_id]
                                        link_text = ''.join(node.text for node in hyperlink.iter() if node.text).strip()
                                        if url and link_text:
                                            all_document_links.append((link_text, url))
                            except Exception:
                                pass

                        if all_document_links:
                            links_section = '\n\n**Links from document:**\n' + '\n'.join([f"- {text}: {url}" for text, url in all_document_links])
                            parts.append(links_section)
                    except Exception:
                        # Ignore link extraction failures
                        pass

                    # Tables
                    try:
                        for table in doc.tables:
                            for row in table.rows:
                                for cell in row.cells:
                                    cell_text = cell.text.strip()
                                    if cell_text:
                                        parts.append(cell_text)
                    except Exception:
                        pass

                    return '\n'.join(parts)
                except Exception:
                    return f"[Simulated DOCX parsing for {filename} - python-docx missing or failed]"

            # PPTX slides
            elif lower.endswith('.pptx'):
                try:
                    from pptx import Presentation
                    prs = Presentation(io.BytesIO(content))
                    text = ''
                    for slide_num, slide in enumerate(prs.slides, 1):
                        text += f"\n--- Slide {slide_num} ---\n"
                        for shape in slide.shapes:
                            shape_text = getattr(shape, "text", None)
                            if shape_text:
                                text += shape_text + "\n"
                    return text
                except Exception:
                    return f"[Simulated PPTX parsing for {filename} - python-pptx missing or failed]"

            # URL shortcut files (.url)
            elif lower.endswith('.url'):
                try:
                    text = content.decode('utf-8', errors='replace')
                    url = None
                    for line in text.splitlines():
                        line = line.strip()
                        if line.upper().startswith('URL='):
                            url = line[4:]
                            break

                    if url:
                        import urllib.parse
                        parsed = urllib.parse.urlparse(str(url))
                        domain = str(parsed.netloc) or 'Unknown domain'
                        filename_without_ext = filename.rsplit('.', 1)[0].replace('_', ' ')

                        content_out = f"**{filename_without_ext}**\n\n"
                        content_out += f"Type: Internet Link / Shortcut\n"
                        content_out += f"Description: This is a link to {filename_without_ext}.\n"

                        if 'powerbi.com' in domain:
                            content_out += "This is a Power BI report link.\n"
                        elif 'sharepoint.com' in domain:
                            content_out += "This is a SharePoint link.\n"

                        content_out += f"\nLink: [{filename_without_ext}]({url})\n"
                        content_out += f"\nKeywords: {filename_without_ext}, link, shortcut, URL\n"

                        return content_out
                    else:
                        return f"Internet shortcut: {filename}"

                except Exception as e:
                    self.logger.exception('Error parsing URL file: %s', e)
                    return f"URL shortcut: {filename}"

            # Fallback
            else:
                return f"[Simulated parsing for {filename} - demo mode]\nContent not extracted."

        except Exception as e:
            self.logger.exception('Error parsing content: %s', e)
            return f"[Error parsing {filename}, using fallback text.]"

# Create a single global parser instance for convenience
parser_service = ParserService()

async def parse_content(filename: str, content: bytes) -> str:
    """Compatibility wrapper that delegates to the shared ParserService.
    Runs the blocking parser in a separate thread to avoid blocking the event loop.
    """
    return await asyncio.to_thread(parser_service.parse_content, filename, content)
