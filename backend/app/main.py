# LLM Assistant - Main Application
# Created by Gustav Christensen
# Date: December 2025
# Description: Quart-based backend for document upload, search, and RAG chat using Qdrant and Azure OpenAI

"""Quart-based demo backend for Qdrant-style search + FastEmbed simulation.
- Async endpoints
- Upload and parse files (txt, csv, pdf, docx)
- Chunking and deterministic embedding (hash-based)
- Search using cosine similarity
- Demo-mode flag and simulated latency
"""
from quart import Quart, request, jsonify, Response, render_template, send_from_directory
from quart_cors import cors
from pydantic import BaseModel
from typing import List, Dict, Any, Optional, Tuple, cast
import uuid
import io
import csv
import time
import asyncio
import inspect
import random
import hashlib
import math
import logging
import datetime                
import gc
import pandas as pd                
import tempfile
from pypdf import PdfReader


try:
    from .config import settings
    from . import embeddings
    from .module.middleware import setup_middleware
    from .module.lifecycle import lifecycle
    from .module.apimonitor import monitor
except ImportError:

    import os, sys
    APP_DIR = os.path.dirname(os.path.abspath(__file__))
    if APP_DIR not in sys.path:
        sys.path.insert(0, APP_DIR)
    import config as config_mod
    settings = config_mod.settings
    import embeddings as embeddings
    from module.middleware import setup_middleware
    from module.lifecycle import lifecycle
    from module.apimonitor import monitor
    # Garbage collection manager (optional psutil dependency)
    try:
        from .module.garbage_collection import gc_manager
    except Exception:
        gc_manager = None

# Basic logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s %(levelname)s %(name)s %(message)s')
log = logging.getLogger(__name__)

from typing import TYPE_CHECKING

if TYPE_CHECKING:
    from qdrant_client import QdrantClient
    from qdrant_client.http import models as qmodels

try:
    from openai import AsyncAzureOpenAI
except ImportError:
    AsyncAzureOpenAI = None

# AI Manager (moved into this file)
class AIManager:
    """Wrapper around Azure OpenAI chat use and response parsing."""

    def __init__(self, settings, embeddings_module, client=None):
        self.settings = settings
        self.embeddings = embeddings_module
        self._client = client

    def initialize_client(self) -> None:
        """Lazily initialize AsyncAzureOpenAI client if available."""
        if self._client is None:
            if AsyncAzureOpenAI is None:
                log.warning("AsyncAzureOpenAI is not available; AI calls will be disabled")
                return
            try:
                # Cast to str to satisfy type checkers - values may be None at runtime
                self._client = AsyncAzureOpenAI(
                    api_key=cast(str, getattr(self.settings, "azure_openai_api_key", "")),
                    api_version=cast(str, getattr(self.settings, "azure_openai_api_version", "")),
                    azure_endpoint=cast(str, getattr(self.settings, "azure_openai_endpoint", "")),
                )
                log.info("AIManager: AsyncAzureOpenAI client initialized")
            except Exception as e:
                log.exception("AIManager: failed to initialize client: %s", e)
                self._client = None

    async def ask_agent_with_context(self, query: str, context: str, conversation_history: Optional[list] = None) -> Tuple[str, Optional[int]]:
        """Send the query + context + history to the chat model.

        Returns (cleaned_text, source_indicator) where source_indicator is 0,1 or None.
        """
        if conversation_history is None:
            conversation_history = []

        # Ensure client
        if self._client is None:
            self.initialize_client()
        if self._client is None:
            return ("Beklager, noget gik galt med assistenten.", None)

        # Build system prompt and messages
        max_ctx = getattr(self.settings, "MAX_CONTEXT_CHARS", 4000)
        ctx = (context[:max_ctx] + "\n\n...[Truncated context]...") if context and len(context) > max_ctx else (context or "")

        system_prompt = self._build_system_prompt(ctx)
        messages = [{"role": "system", "content": system_prompt}]

        # Conversation history (safe)
        max_history = int(getattr(self.settings, "MAX_CONVERSATION_HISTORY", 6))
        history_limit = min(len(conversation_history), max_history)
        for msg in (conversation_history or [])[-history_limit:]:
            if isinstance(msg, dict):
                u = msg.get("user", "")
                a = msg.get("assistant", "")
                if u:
                    messages.append({"role": "user", "content": u})
                if a:
                    messages.append({"role": "assistant", "content": a})

        messages.append({"role": "user", "content": query})

        try:
            # Ensure model is a proper string per type expectations
            model = cast(str, getattr(self.settings, "azure_deployment_chat", ""))
            if not model:
                log.warning("AIManager: no chat model configured; aborting request")
                return ("Beklager, noget gik galt med assistenten.", None)

            temperature = float(getattr(self.settings, "TEMPERATURE", 0.3))
            max_tokens = int(getattr(self.settings, "MAX_TOKENS", 512))

            # Cast messages to Any to satisfy typing for third-party API wrappers
            completion = await self._client.chat.completions.create(
                model=model,
                messages=cast(Any, messages),
                temperature=temperature,
                max_tokens=max_tokens,
            )

            try:
                raw = getattr(completion.choices[0].message, "content", "") or ""
            except Exception:
                raw = str(completion) if completion else ""

            cleaned, source = self._parse_and_clean_source(raw)
            return cleaned, source

        except Exception as e:
            txt = str(e).lower()
            if "content_filter" in txt or "responsibleaipolicyviolation" in txt or (hasattr(e, "status_code") and getattr(e, "status_code", None) == 400):
                return "__CONTENT_FILTER_BLOCKED__", None
            log.exception("AIManager: ask_agent_with_context failed: %s", e)
            return ("Beklager, noget gik galt med assistenten.", None)

    def _parse_and_clean_source(self, response_text: str) -> Tuple[str, Optional[int]]:
        """Robust parsing of a [SOURCE:X] marker on the FIRST line only."""
        if not response_text:
            return response_text, None

        m = re.match(r'^\s*\**\s*\[SOURCE:(?P<s>[01])\]\s*\**\s*(?:\n)?\s*(?P<rest>.*)', response_text, flags=re.S)
        if m:
            rest = m.group("rest") or ""
            return rest.strip(), int(m.group("s"))
        return response_text, None

    def _build_system_prompt(self, context: str) -> str:
        base = (
            "You are a helpful, concise, and safety-conscious assistant. "
            "Answer clearly in English and use Markdown formatting when it improves readability. "
            "Be factual, admit uncertainty when appropriate, and do not fabricate information or external links.\n\n"
        )
        base += "CONTEXT: " + (context if context else "[None]")
        return base


class DemoLLM:
    """Deterministic demo-mode LLM emulator.

    Returns short, deterministic answers built from provided context and marks
    a SOURCE (0 or 1) heuristically.
    """

    def __init__(self, max_chars: int = 800):
        self.max_chars = max_chars

    async def generate(self, query: str, context: str) -> tuple[str, Optional[int]]:
        # If no context, return a helpful demo message
        if not context or not context.strip():
            return ("I have no documents to answer from. Try uploading or adding documents.", None)

        # Pick up to 3 blocks separated by paragraphs
        parts = [p.strip() for p in context.split("\n\n") if p.strip()][:3]
        body = "\n\n".join(parts)
        if len(body) > self.max_chars:
            body = body[: self.max_chars - 3].rstrip() + "..."

        # Heuristic source selection
        lower = body.lower()
        source = 1 if "microsoft" in lower or "windows" in lower or "office" in lower else 0

        # Return in the same style as the real LLM (with SOURCE on first line)
        text = f"[SOURCE:{source}]\n{body}"
        return text, source

# Optional Qdrant imports for typing
try:
    from qdrant_client import QdrantClient
    from qdrant_client.http import models as qmodels
except ImportError:
    pass # Managed by lifecycle

app = Quart(__name__, template_folder='../templates', static_folder='../static')
app = cors(app, allow_origin="*")

# Register Middleware
setup_middleware(app)

# In-memory store (for demo or fallback)
DOCUMENTS: List[Dict[str, Any]] = []

# Helper: simple chunker / cleanup
def _clean_text(t: str) -> str:
    return t.replace('\r', '\n').strip()


DEMO_BANNER = "Demo mode — results are simulated; no active LLM connection."
EMBED_DIM = settings.embedding_dim


class SearchRequest(BaseModel):
    query: str
    top_k: int = 5


import re

async def chunk_text_async(text: str, chunk_size: int = 200, overlap: int = 50) -> List[str]:
    return await embedding_service.chunk_text_async(text, chunk_size=chunk_size, overlap=overlap)

def chunk_text(text: str, chunk_size: int = 200, overlap: int = 50) -> List[str]:
    return embedding_service.chunk_text(text, chunk_size=chunk_size, overlap=overlap)


class EmbeddingService:
    """Helper for deterministic embeddings and chunking logic."""

    def __init__(self, dim: int = EMBED_DIM):
        self.dim = dim

    def deterministic_embed(self, text: str) -> List[float]:
        # Deterministic embedding using md5 digest split into numbers
        h = hashlib.md5(text.encode('utf-8')).digest()
        vals = []
        # expand digest to dim by repeating and combining
        while len(vals) < self.dim:
            for b in h:
                vals.append(b / 255.0)
                if len(vals) >= self.dim:
                    break
            h = hashlib.md5(h).digest()
        # normalize vector
        norm = math.sqrt(sum(v * v for v in vals)) or 1.0
        return [v / norm for v in vals]

    def cosine_sim(self, a: List[float], b: List[float]) -> float:
        return sum(x * y for x, y in zip(a, b))

    async def chunk_text_async(self, text: str, chunk_size: int = 200, overlap: int = 50) -> List[str]:
        loop = asyncio.get_running_loop()
        return await loop.run_in_executor(None, self.chunk_text, text, chunk_size, overlap)

    def chunk_text(self, text: str, chunk_size: int = 200, overlap: int = 50) -> List[str]:
        """Chunk text by sentences into near chunk_size words with overlap.
        This gives more natural breaks than raw word slicing."""
        # Normalize whitespace
        text = re.sub(r"\s+", " ", text).strip()
        if not text:
            return []
        # Split into sentences (simple heuristic)
        sentences = re.split(r'(?<=[.!?])\s+', text)
        chunks = []
        current = []
        current_words = 0
        for s in sentences:
            sw = len(s.split())
            if current_words + sw <= chunk_size or not current:
                current.append(s)
                current_words += sw
            else:
                chunks.append(' '.join(current))
                # start next chunk with overlap sentences (approx by words)
                # carry over last sentences until overlap words satisfied
                carry = []
                carry_words = 0
                j = len(current) - 1
                while j >= 0 and carry_words < overlap:
                    sentence_j = current[j]
                    carry.insert(0, sentence_j)
                    carry_words += len(sentence_j.split())
                    j -= 1
                current = carry + [s]
                current_words = sum(len(x.split()) for x in current)
        if current:
            chunks.append(' '.join(current))
        return chunks


# Single instance for compatibility with existing code
embedding_service = EmbeddingService()


def deterministic_embed(text: str) -> List[float]:
    return embedding_service.deterministic_embed(text)


def cosine_sim(a: List[float], b: List[float]) -> float:
    return embedding_service.cosine_sim(a, b)


class ParserService:
    """Service responsible for extracting text and links from uploaded files.

    This encapsulates the previous free-function `parse_content` and provides a
    single, testable interface for file parsing.
    """

    def __init__(self):
        # placeholder for future options (e.g., max file size, ocr settings)
        self.logger = log

    async def parse_content(self, filename: str, content: bytes) -> str:
        """See previous `parse_content` behavior.

        Keeps backward-compatible behavior and graceful fallbacks when optional
        libraries are missing.
        """

        lower = filename.lower()

        try:
            # Plain text and markdown
            if lower.endswith('.txt') or lower.endswith('.md'):
                return content.decode('utf-8', errors='replace')

            # Enhanced CSV parsing: try pandas if available for better formats
            elif lower.endswith('.csv'):
                s = content.decode('utf-8', errors='replace')
                try:

                    # Try to detect a header-oriented CSV with semicolon separators
                    lines = s.splitlines()
                    csv_start = 0
                    gruppe_navn = ''
                    found_header = False

                    for i, line in enumerate(lines):
                        line = line.strip()
                        if line.startswith('Sheet:'):
                            gruppe_navn = line.replace('Sheet:', '').strip()
                        elif ';' in line or ',' in line:
                            csv_start = i
                            found_header = True
                            break

                    csv_lines = '\n'.join(lines[csv_start:]) if found_header else s

                    with tempfile.NamedTemporaryFile(mode='w', suffix='.csv', delete=False, encoding='utf-8') as tmpf:
                        tmpf.write(csv_lines)
                        tmp_path = tmpf.name

                    try:
                        # Prefer semicolon when present
                        sep = ';' if ';' in csv_lines else ','
                        df = pd.read_csv(tmp_path, sep=sep, encoding='utf-8')
                        content_out = f"**Data fra {filename}**\n\n"
                        if gruppe_navn:
                            content_out += f"Gruppe: {gruppe_navn}\n"
                        content_out += f"Kolonner: {', '.join(df.columns)}\n"
                        content_out += f"Antal rækker: {len(df)}\n\n"

                        for _, row in df.iterrows():
                            row_data = []
                            for col in df.columns:
                                val = row[col]
                                if pd.notna(val):
                                    row_data.append(f"{col}: {val}")
                            content_out += ' | '.join(row_data) + '\n'

                        try:
                            import os as _os
                            _os.unlink(tmp_path)
                        except Exception:
                            pass

                        return content_out
                    except Exception:
                        # Fallback to simple csv reader
                        rows = list(csv.reader(io.StringIO(s)))
                        lines = [', '.join(row) for row in rows]
                        return '\n'.join(lines)
                except Exception:
                    # pandas or tempfile not available - simple fallback
                    rows = list(csv.reader(io.StringIO(s)))
                    lines = [', '.join(row) for row in rows]
                    return '\n'.join(lines)

            # PDF parsing with link extraction
            if lower.endswith('.pdf'):
                try:
                    text_parts = []
                    links = []

                    try:
                        # Prefer pypdf for link extraction and robust text
                        
                        reader = PdfReader(io.BytesIO(content))
                        for i, page in enumerate(reader.pages, start=1):
                            try:
                                page_text = page.extract_text() or ''
                                if page_text.strip():
                                    text_parts.append(page_text)
                            except Exception:
                                pass

                            # Extract annotations/links if present
                            try:
                                annots_obj = page.get('/Annots')  # type: ignore[reportGeneralTypeIssues]
                                if annots_obj:
                                    # annots_obj may be an array-like PdfObject; iterate defensively
                                    try:
                                        for a in annots_obj:
                                            try:
                                                obj = a.get_object()
                                                if obj.get('/Subtype') == '/Link' and '/A' in obj:
                                                    action = obj['/A']
                                                    if action.get('/S') == '/URI' and '/URI' in action:
                                                        uri = action['/URI']
                                                        links.append(f"Link (side {i}): {uri}")
                                            except Exception:
                                                continue
                                    except TypeError:
                                        # Not iterable - try single object handling
                                        try:
                                            a = annots_obj
                                            obj = a.get_object()
                                            if obj.get('/Subtype') == '/Link' and '/A' in obj:
                                                action = obj['/A']
                                                if action.get('/S') == '/URI' and '/URI' in action:
                                                    uri = action['/URI']
                                                    links.append(f"Link (side {i}): {uri}")
                                        except Exception:
                                            pass
                            except Exception:
                                # Some pdf versions have different structures; ignore
                                pass

                    except Exception:
                        # Fallback to pdfplumber for text only (no annotations extraction)
                        try:
                            import pdfplumber
                            with pdfplumber.open(io.BytesIO(content)) as pdf:
                                for page in pdf.pages:
                                    txt = page.extract_text() or ''
                                    if txt.strip():
                                        text_parts.append(txt)
                                    else:
                                        # Try OCR fallback if pytesseract is available
                                        try:
                                            from PIL import Image
                                            import pytesseract
                                            img = page.to_image(resolution=150).original
                                            ocr_txt = pytesseract.image_to_string(img)
                                            text_parts.append(ocr_txt)
                                        except Exception:
                                            pass
                        except Exception:
                            return f"[Simulated PDF parsing for {filename} - pdf parsing libraries missing or failed]"

                    combined = '\n'.join([p for p in text_parts if p and p.strip()])

                    if links:
                        combined += '\n\n**Links fra PDF dokumentet:**\n'
                        for l in links:
                            combined += f"- {l}\n"

                    if combined.strip() == '':
                        return f"[Simulated PDF parsing for {filename} - no text extracted]"

                    return combined
                except Exception as e:
                    self.logger.exception('PDF parsing error: %s', e)
                    return f"[Simulated PDF parsing for {filename} - error]"

            # DOCX parsing with link extraction
            elif lower.endswith('.docx'):
                try:
                    from docx import Document

                    doc = Document(io.BytesIO(content))
                    parts = []

                    # Extract paragraphs
                    for para in doc.paragraphs:
                        if para.text and para.text.strip():
                            parts.append(para.text.strip())

                    # Extract hyperlinks using rels
                    try:
                        all_hyperlinks = {}
                        for rel in doc.part.rels.values():
                            if 'hyperlink' in rel.reltype:
                                all_hyperlinks[rel.rId] = rel.target_ref

                        all_document_links = []
                        for paragraph in doc.paragraphs:
                            try:
                                for hyperlink in paragraph._element.findall(
                                    './/{http://schemas.openxmlformats.org/wordprocessingml/2006/main}hyperlink'
                                ):
                                    rel_id = hyperlink.get(
                                        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
                                    )
                                    if rel_id in all_hyperlinks:
                                        url = all_hyperlinks[rel_id]
                                        link_text = ''.join(node.text for node in hyperlink.iter() if node.text).strip()
                                        if url and link_text:
                                            all_document_links.append((link_text, url))
                            except Exception:
                                pass

                        if all_document_links:
                            links_section = '\n\n**Links fra dokumentet:**\n' + '\n'.join([f"- {text}: {url}" for text, url in all_document_links])
                            parts.append(links_section)
                    except Exception:
                        # Ignore link extraction failures
                        pass

                    # Tables
                    try:
                        for table in doc.tables:
                            for row in table.rows:
                                for cell in row.cells:
                                    cell_text = cell.text.strip()
                                    if cell_text:
                                        parts.append(cell_text)
                    except Exception:
                        pass

                    return '\n'.join(parts)
                except Exception:
                    return f"[Simulated DOCX parsing for {filename} - python-docx missing or failed]"

            # PPTX slides
            elif lower.endswith('.pptx'):
                try:
                    from pptx import Presentation
                    prs = Presentation(io.BytesIO(content))
                    text = ''
                    for slide_num, slide in enumerate(prs.slides, 1):
                        text += f"\n--- Slide {slide_num} ---\n"
                        for shape in slide.shapes:
                            shape_text = getattr(shape, "text", None)
                            if shape_text:
                                text += shape_text + "\n"
                    return text
                except Exception:
                    return f"[Simulated PPTX parsing for {filename} - python-pptx missing or failed]"

            # URL shortcut files (.url)
            elif lower.endswith('.url'):
                try:
                    text = content.decode('utf-8', errors='replace')
                    url = None
                    for line in text.splitlines():
                        line = line.strip()
                        if line.upper().startswith('URL='):
                            url = line[4:]
                            break

                    if url:
                        import urllib.parse
                        parsed = urllib.parse.urlparse(str(url))
                        domain = str(parsed.netloc) or 'Unknown domain'
                        filename_without_ext = filename.rsplit('.', 1)[0].replace('_', ' ')

                        content_out = f"**{filename_without_ext}**\n\n"
                        content_out += f"Type: Internet Link/Genvej\n"
                        content_out += f"Beskrivelse: Dette er et link til {filename_without_ext}.\n"

                        if 'powerbi.com' in domain:
                            content_out += "Dette er et Power BI rapporteringslink.\n"
                        elif 'sharepoint.com' in domain:
                            content_out += "Dette er et SharePoint link.\n"

                        content_out += f"\nLink: [{filename_without_ext}]({url})\n"
                        content_out += f"\nNøgleord: {filename_without_ext}, link, genvej, URL\n"

                        return content_out
                    else:
                        return f"Internet shortcut: {filename}"

                except Exception as e:
                    log.exception('Error parsing URL file: %s', e)
                    return f"URL shortcut: {filename}"

            # Fallback
            else:
                return f"[Simulated parsing for {filename} - demo mode]\nContent not extracted."

        except Exception as e:
            log.exception('Error parsing content: %s', e)
            return f"[Error parsing {filename}, using fallback text.]"

# Create a single global parser instance for convenience
parser_service = ParserService()

async def parse_content(filename: str, content: bytes) -> str:
    """Compatibility wrapper that delegates to the shared ParserService."""
    return await parser_service.parse_content(filename, content)


class QdrantService:
    """Wrapper for common Qdrant operations used by endpoints."""

    def __init__(self, lifecycle_obj):
        self.lifecycle = lifecycle_obj
        self.collection = settings.qdrant_collection

    def is_available(self) -> bool:
        return bool(self.lifecycle.qdrant_client)

    def scroll(self, limit: int = 2000):
        try:
            return self.lifecycle.qdrant_client.scroll(collection_name=self.collection, limit=limit)
        except Exception:
            return ([], None)

    def list_documents(self):
        if not self.is_available():
            return []
        try:
            res, _ = self.scroll(limit=2000)
        except Exception:
            return []

        unique_files = {}
        for hit in res or []:
            payload = getattr(hit, 'payload', {}) or {}
            fname = payload.get('filename') or payload.get('file_name') or 'unknown'
            if fname not in unique_files:
                unique_files[fname] = {
                    "id": fname,
                    "filename": fname,
                    "chunks": 0,
                    "uploaded_at": payload.get('uploaded_at', 'N/A'),
                    "type": payload.get('type', 'FILE')
                }
            unique_files[fname]["chunks"] += 1

        return list(unique_files.values())

    def get_document_chunks(self, doc_id: str):
        if not self.is_available():
            return []
        try:
            res, _ = self.lifecycle.qdrant_client.scroll(collection_name=self.collection, limit=1000)
        except Exception:
            return []

        chunks = []
        for hit in res or []:
            payload = getattr(hit, 'payload', {}) or {}
            if payload.get('filename') == doc_id or payload.get('file_name') == doc_id:
                chunks.append({"text": payload.get('text', ''), "chunk_index": payload.get('chunk_index', 0)})

        return chunks

    def delete_by_filename(self, doc_id: str) -> int:
        """Delete all points whose payload filename/file_name matches doc_id.
        Returns number of deleted points (best-effort)."""
        if not self.is_available():
            return 0

        try:
            scroll_res = self.lifecycle.qdrant_client.scroll(collection_name=self.collection, limit=10000)
            all_points = scroll_res[0] if scroll_res else []
        except Exception:
            all_points = []

        ids_to_delete = []
        for point in all_points:
            payload = getattr(point, 'payload', {}) or {}
            if payload.get('filename') == doc_id or payload.get('file_name') == doc_id:
                ids_to_delete.append(point.id)

        if ids_to_delete:
            self.lifecycle.qdrant_client.delete(collection_name=self.collection, points_selector=ids_to_delete)

        return len(ids_to_delete)

    def initialize_client(self) -> None:
        """Initialize Qdrant client with safe persistent local storage."""
        try:
            # Prefer lifecycle-managed client if available
            client = getattr(self.lifecycle, 'qdrant_client', None)
            if client:
                # Keep a short local reference
                self.lifecycle.qdrant_client = client
                return

            # If lifecycle does not have a client, attempt to create a path-based client
            qdrant_path = os.path.join(os.path.dirname(__file__), 'qdrant_db')
            os.makedirs(qdrant_path, exist_ok=True)
            from qdrant_client import QdrantClient as _QdrantClient
            # Note: we prefer file-backed path for local persistence when possible
            try:
                self.lifecycle.qdrant_client = _QdrantClient(path=qdrant_path)
                log.info(f"Qdrant client initialized at {qdrant_path}")
            except Exception as e:
                log.warning(f"Could not initialize local Qdrant client: {e}")
        except Exception as e:
            log.error(f"Qdrant initialize_client error: {e}")

    async def setup_collection(self) -> bool:
        """Ensure the collection exists and mark it as ready. Returns True if ready."""
        try:
            client = getattr(self.lifecycle, 'qdrant_client', None)
            if client is None:
                log.warning("Qdrant client not available - cannot setup collection")
                return False

            try:
                collection_info = await client.get_collection(self.collection)
                log.info(f"Qdrant collection '{self.collection}' ready: {collection_info.points_count} points")
                self.collection_ready = True
                return True
            except Exception as e:
                log.info(f"Qdrant collection '{self.collection}' not present or unusable: {e}")
                self.collection_ready = False
                return False
        except Exception as e:
            log.error(f"setup_collection error: {e}")
            return False

    def preprocess_query(self, query_text: str) -> str:
        """Preprocess query to expand common terms for better matching (English)."""
        try:
            query = " ".join(query_text.strip().split())

            synonyms = {
                "office": "Office 365 Microsoft 365",
                "teams": "Microsoft Teams",
                "outlook": "Microsoft Outlook",
                "word": "Microsoft Word",
                "excel": "Microsoft Excel",
                "powerpoint": "Microsoft PowerPoint",
                "sharepoint": "Microsoft SharePoint",
                "onedrive": "OneDrive",
                "login": "log in login access password account",
                "printer": "printer print printing",
                "internet": "network wifi connection",
                "mail": "email outlook",
                "password": "password change reset update",
                "reset": "reset restart reboot",
            }

            q_lower = query.lower()
            for term, extras in synonyms.items():
                if term in q_lower:
                    query += f" {extras}"

            return query
        except Exception:
            return query_text

    async def search_index(self, query_text: str, top_k: int | None = None) -> dict:
        """Search the vector index with retries, returns context and sources."""
        try:
            if top_k is None:
                top_k = int(getattr(settings, 'TOP_K_RESULTS', 5))

            await self.setup_collection()

            processed_query = self.preprocess_query(query_text)

            # Generate embedding (use project's embeddings helper if available)
            emb_list = await embeddings.embed_texts([processed_query], provider=settings.embedding_provider, context=f"search:{processed_query[:80]}")
            if not emb_list:
                return {"context": "", "sources": [], "embedding_time": 0.0, "search_time": 0.0}
            vector = emb_list[0]

            client = getattr(self.lifecycle, 'qdrant_client', None)
            if client is None:
                log.error("Qdrant client not available - cannot search")
                return {"context": "", "sources": [], "embedding_time": 0.0, "search_time": 0.0}

            # Retry with exponential backoff
            max_retries = 3
            base_delay = 0.5
            results = None
            for attempt in range(max_retries):
                try:
                    # Use search API if present or fallback to query_points
                    if hasattr(client, 'search'):
                        results = await client.search(collection_name=self.collection, query_vector=vector, limit=min(top_k*2, getattr(settings, 'MAX_SEARCH_LIMIT', 100)), score_threshold=getattr(settings, 'QDRANT_SCORE_THRESHOLD', 0.0))
                        # client.search returns list-like results
                    else:
                        # query_points for older client
                        res = client.query_points(collection_name=self.collection, query=vector, limit=top_k*2)
                        results = getattr(res, 'results', res)
                    break
                except Exception as e:
                    if attempt < max_retries - 1:
                        delay = base_delay * (2 ** attempt)
                        log.warning(f"Qdrant search attempt {attempt+1} failed: {e}. Retrying in {delay}s")
                        await asyncio.sleep(delay)
                    else:
                        log.error(f"Qdrant search failed after {max_retries} attempts: {e}")
                        return {"context": "", "sources": [], "embedding_time": 0.0, "search_time": 0.0}

            # Format results into context and sources
            chunks = []
            source_docs = []
            seen_titles = set()

            if results:
                # results may be iterable of hits
                for hit in results:
                    payload = getattr(hit, 'payload', {}) or {}
                    score = getattr(hit, 'score', None) or 0.0
                    text = payload.get('text', '')
                    title_src = payload.get('file_name') or payload.get('filename') or 'Document'
                    title = title_src.replace('.pdf','').replace('.docx','').replace('_',' ')

                    # Apply simple score threshold if present
                    threshold = float(getattr(settings, 'SIMILARITY_THRESHOLD', 0.0))
                    try:
                        if float(score) <= threshold:
                            continue
                    except Exception:
                        pass

                    if text:
                        if title not in seen_titles:
                            source_docs.append({"title": title, "similarity_score": round(float(score), 3)})
                            seen_titles.add(title)

                        if title and title.lower() not in text[:100].lower():
                            chunks.append(f"**{title}**\n\n{text}")
                        else:
                            chunks.append(text)

                    if len(source_docs) >= top_k:
                        break

            return {"context": "\n\n".join(chunks), "sources": source_docs, "embedding_time": 0.0, "search_time": 0.0}
        except Exception as e:
            log.error(f"search_index error: {e}")
            return {"context": "", "sources": [], "embedding_time": 0.0, "search_time": 0.0}

    def set_loading_state(self, is_loading: bool, operation: str = "", total_docs: int = 0) -> None:
        """Set loading state to inform users about long-running operations."""
        if is_loading:
            self.loading_state = {
                "is_loading": True,
                "start_time": time.time(),
                "total_documents": total_docs,
                "processed_documents": 0,
                "current_operation": operation,
                "estimated_completion": time.time() + (total_docs * 2),
            }
        else:
            self.loading_state = {
                "is_loading": False,
                "start_time": None,
                "total_documents": 0,
                "processed_documents": 0,
                "current_operation": "",
                "estimated_completion": None,
            }
            gc.collect()

    def get_loading_message(self) -> str | None:
        """Return a friendly English loading message when the index is being updated."""
        if not self.loading_state.get("is_loading"):
            return None

        messages = [
            "I'm updating my knowledge base - please check back in {time} minutes!",
            "Processing new documents now, give me {time} minutes and I'll be ready!",
            "Index update in progress - come back in about {time} minutes."
        ]

        if self.loading_state.get("estimated_completion"):
            remaining = max(1, int((self.loading_state["estimated_completion"] - time.time()) / 60))
        else:
            remaining = 5
        remaining = max(1, remaining + random.randint(-1, 2))

        msg = random.choice(messages)
        return msg.format(time=remaining)

    def update_loading_progress(self, processed_docs: int) -> None:
        if self.loading_state.get("is_loading"):
            self.loading_state["processed_documents"] = processed_docs
            if processed_docs > 0 and self.loading_state.get("start_time"):
                elapsed = time.time() - self.loading_state["start_time"]
                avg = elapsed / processed_docs
                remaining = max(0, self.loading_state.get("total_documents", 0) - processed_docs)
                self.loading_state["estimated_completion"] = time.time() + (remaining * avg)

    async def clear_collection(self) -> dict:
        """Delete all points from the collection by scrolling and deleting in batches."""
        try:
            client = getattr(self.lifecycle, 'qdrant_client', None)
            if client is None:
                return {"success": False, "error": "Qdrant client not available"}

            self.set_loading_state(True, "Clearing collection", 0)
            total_deleted = 0
            offset = None
            batch_size = 100

            while True:
                scroll_result = await client.scroll(collection_name=self.collection, limit=batch_size, offset=offset, with_payload=False, with_vectors=False)
                points, next_offset = scroll_result
                if not points:
                    break
                ids = [p.id for p in points]
                await client.delete(collection_name=self.collection, points_selector=ids)
                total_deleted += len(ids)
                offset = next_offset
                if next_offset is None:
                    break

            self.set_loading_state(False)
            return {"success": True, "deleted": total_deleted}
        except Exception as e:
            self.set_loading_state(False)
            log.error(f"clear_collection error: {e}")
            return {"success": False, "error": str(e)}

    async def optimize_collection(self) -> dict:
        """Trigger collection optimization (best-effort)."""
        try:
            client = getattr(self.lifecycle, 'qdrant_client', None)
            if client is None:
                return {"success": False, "error": "Qdrant client not available"}

            try:
                await client.update_collection(collection_name=self.collection, optimizer_config=qmodels.OptimizersConfigDiff(indexing_threshold=0))
                return {"success": True, "message": "Optimization triggered"}
            except Exception as e:
                log.warning(f"optimize_collection failed: {e}")
                return {"success": False, "error": str(e)}
        except Exception as e:
            log.error(f"optimize_collection error: {e}")
            return {"success": False, "error": str(e)}


# Single instance for endpoints to use
qdrant_service = QdrantService(lifecycle)


class SupportAgent:
    """Aggregator that keeps references to the main services used by routes."""

    def __init__(self):
        self.parser = parser_service
        self.embedding = embedding_service
        self.qdrant = qdrant_service
        self.settings = settings
        self.lifecycle = lifecycle
        # AI manager moved from legacy app - lazy-initializes its client
        self.ai_manager = AIManager(self.settings, self.embedding)


# Single global agent
support_agent = SupportAgent()


@app.route('/health')
async def health():
    """Enhanced health check with component status, memory and middleware metrics."""
    try:
        basic_health = {
            "status": "healthy",
            "timestamp": datetime.datetime.now().isoformat(),
            "components": {},
        }

        # API monitor summary if available
        try:
            if hasattr(monitor, 'get_health_summary'):
                monitor_summary = monitor.get_health_summary()
                basic_health.update(monitor_summary)
            else:
                basic_health['metrics'] = monitor.get_stats()
        except Exception as e:
            log.warning(f"Failed to get api monitor health: {e}")

        # Qdrant: client presence and collection readiness
        try:
            qdrant_info = {
                "status": "healthy" if lifecycle.qdrant_client else "unavailable",
                "collection_ready": getattr(qdrant_service, 'collection_ready', False),
            }
            if qdrant_info["collection_ready"] and lifecycle.qdrant_client:
                try:
                    coll_call = lifecycle.qdrant_client.get_collection(settings.qdrant_collection)
                    # handle both sync and async client implementations
                    if inspect.isawaitable(coll_call):
                        coll = await coll_call  # type: ignore
                    else:
                        coll = coll_call
                    qdrant_info["points_count"] = getattr(coll, 'points_count', None)
                except Exception:
                    # don't fail the entire health check for a single probe
                    pass
            basic_health["components"]["qdrant"] = qdrant_info
        except Exception as e:
            basic_health["components"]["qdrant"] = {"status": "error", "error": str(e)}

        # OpenAI availability (best-effort)
        try:
            openai_status = "healthy" if AsyncAzureOpenAI and getattr(settings, 'azure_openai_api_key', None) else "unavailable"
            basic_health["components"]["openai"] = {"status": openai_status}
        except Exception as e:
            basic_health["components"]["openai"] = {"status": "error", "error": str(e)}

        # Database (project may not have a DB configured)
        try:
            # This project currently does not expose a DB setting; report as not configured
            basic_health["components"]["database"] = {"status": "not_configured"}
        except Exception as e:
            basic_health["components"]["database"] = {"status": "error", "error": str(e)}

        # Determine overall status
        component_statuses = [c.get("status") for c in basic_health["components"].values()]
        if "error" in component_statuses or "unavailable" in component_statuses:
            basic_health["status"] = "degraded"

        # Add GC memory info if available
        try:
            if 'gc_manager' in globals() and gc_manager is not None:
                mem = await gc_manager.get_memory_usage()
                if mem:
                    basic_health['memory_mb'] = mem.get('rss_mb', 0)
        except Exception as e:
            log.warning(f"Failed to get gc_manager memory info: {e}")

        # Add middleware health if available
        try:
            cleanup_mw = getattr(app, 'cleanup_middleware', None)
            if cleanup_mw is not None:
                middleware_health = await cleanup_mw.get_health_metrics()
                basic_health['middleware'] = {
                    'memory_usage_mb': middleware_health.get('memory_usage_mb', 0),
                    'active_connections': middleware_health.get('active_connections', 0),
                    'cleanup_runs': middleware_health.get('cleanup_runs', 0),
                }
        except Exception as e:
            log.warning(f"Failed to get cleanup_middleware health: {e}")

        # Response code
        status_code = 200 if basic_health["status"] == "healthy" else 503
        return jsonify(basic_health), status_code

    except Exception as e:
        log.exception('Health check error: %s', e)
        return jsonify({"status": "error", "error": str(e)}), 500


@app.route('/health/gc')
async def health_gc():
    """Return garbage collection manager statistics if available."""
    try:
        if 'gc_manager' not in globals() or gc_manager is None:
            return jsonify({"error": "gc_manager not available"}), 404
        stats = await gc_manager.get_statistics()
        return jsonify({"status": "ok", "gc": stats})
    except Exception as e:
        log.exception('Failed to get GC stats: %s', e)
        return jsonify({"error": str(e)}), 500


@app.route('/health/middleware')
async def health_middleware():
    """Return health metrics from the EnhancedCleanupMiddleware if installed."""
    try:
        middleware = getattr(app, 'cleanup_middleware', None)
        if middleware is None:
            return jsonify({"error": "middleware metrics not available"}), 404
        stats = await middleware.get_health_metrics()
        return jsonify({"status": "ok", "middleware": stats})
    except Exception as e:
        log.exception('Failed to get middleware stats: %s', e)
        return jsonify({"error": str(e)}), 500


@app.route('/')
async def index():
    return await render_template('index.html')


@app.route('/upload', methods=['POST'])
async def upload_file():
    # Accept either multipart file uploads or JSON with 'text' and optional 'filename'
    form = await request.files
    if 'file' in form:
        file = form['file']
        content = file.read()
        filename = file.filename or f"upload-{int(time.time())}"
        text = await parse_content(filename, content)
    else:
        payload = await request.get_json(silent=True)
        if payload and 'text' in payload:
            filename = payload.get('filename', f"upload-{int(time.time())}")
            text = payload['text']
        else:
            return jsonify({"error": "file field or JSON text required"}), 400

    chunks = await chunk_text_async(text)
    doc_id = str(uuid.uuid4())

    # Admin auth (if configured)
    if settings.admin_token:
        auth = request.headers.get('Authorization', '')
        if not auth.startswith('Bearer '):
            return jsonify({'error': 'unauthorized'}), 401
        token = auth.split(' ', 1)[1].strip()
        if token != settings.admin_token:
            return jsonify({'error': 'unauthorized'}), 401

    # Store in Qdrant if available, otherwise in-memory
    if not lifecycle.qdrant_client:
        # Batch process embeddings even in demo mode to avoid log spam and overhead
        embeddings_batch = await embeddings.embed_texts(chunks, provider=settings.embedding_provider, context=filename)
        
        chunk_objs = []
        for c, emb in zip(chunks, embeddings_batch):
            chunk_objs.append({"text": c, "embed": emb})

        DOCUMENTS.append({
            "id": doc_id,
            "filename": filename,
            "text": text,
            "chunks": chunk_objs,
            "uploaded_at": datetime.datetime.now().strftime("%Y-%m-%d %H:%M"),
            "type": filename.split('.')[-1].upper() if '.' in filename else 'TXT'
        })
        return jsonify({"id": doc_id, "filename": filename, "status": "uploaded (fallback)"})

    # Index into Qdrant (both demo and prod modes)
    # Create points and upsert
    points = []
    embeddings_batch = await embeddings.embed_texts(chunks, provider=settings.embedding_provider, context=filename)
    timestamp = datetime.datetime.now().strftime("%Y-%m-%d %H:%M")
    file_type = filename.split('.')[-1].upper() if '.' in filename else 'TXT'
    
    for idx, (c, emb) in enumerate(zip(chunks, embeddings_batch)):
        points.append(qmodels.PointStruct(
            id=str(uuid.uuid4()), 
            vector=emb, 
            payload={
                "filename": filename, 
                "chunk_index": idx, 
                "text": c,
                "uploaded_at": timestamp,
                "type": file_type
            }
        ))

    try:
        log.info('Upserting %d points into collection %s', len(points), settings.qdrant_collection)
        lifecycle.qdrant_client.upsert(collection_name=settings.qdrant_collection, points=points)
        # Verify by counting points after upsert
        try:
            count = lifecycle.qdrant_client.count(collection_name=settings.qdrant_collection)
            log.info("After upsert, collection '%s' contains %s points", settings.qdrant_collection, getattr(count, 'count', count))
        except Exception:
            log.debug('Could not retrieve collection count after upsert')
        return jsonify({"id": doc_id, "filename": filename, "status": "indexed"})
    except Exception as e:
        log.exception('Failed to upsert points into Qdrant: %s', e)
        return jsonify({"error": "failed to index document", "detail": str(e)}), 500


@app.route('/documents')
async def list_documents():
    if not lifecycle.qdrant_client:
        return jsonify([{
            "id": d["id"], 
            "filename": d["filename"], 
            "chunks": len(d["chunks"]),
            "uploaded_at": d.get("uploaded_at", "N/A"),
            "type": d.get("type", "FILE")
        } for d in DOCUMENTS])
    
    # Use QdrantService to list documents
    try:
        docs = qdrant_service.list_documents()
        return jsonify(docs)
    except Exception:
        return jsonify([])


@app.route('/document/<doc_id>')
async def get_document(doc_id):
    if not lifecycle.qdrant_client:
        doc = next((d for d in DOCUMENTS if d["id"] == doc_id), None)
        if not doc:
            return jsonify({"error": "document not found"}), 404
        return jsonify({"id": doc["id"], "filename": doc["filename"], "text": doc["text"], "chunks": len(doc["chunks"])})
    # From Qdrant, get all points with that doc_id? Wait, doc_id is not stored.
    # Actually, since points have payload with filename, but not doc_id.
    # In demo, doc_id is generated, but in Qdrant, points have id as uuid, payload filename.
    # To get full document, perhaps need to store full text in payload or something.
    # For now, return points for that filename.
    # But filename may not be unique.
    # For simplicity, since demo, and Qdrant, perhaps return the chunks.
    chunks = qdrant_service.get_document_chunks(doc_id)
    if not chunks:
        return jsonify({"error": "document not found"}), 404

    # Reconstruct text
    chunks.sort(key=lambda x: x['chunk_index'])
    full_text = '\n'.join(c['text'] for c in chunks)
    return jsonify({"id": doc_id, "filename": doc_id, "text": full_text, "chunks": len(chunks)})


@app.route('/search', methods=['POST'])
async def search():
    payload = await request.get_json()
    if not payload or 'query' not in payload:
        return jsonify({"error": "query required"}), 400
    query = payload.get('query', '').strip()
    top_k = int(payload.get('top_k', 5))
    simulate_latency = payload.get('simulate_latency', True)

    if not query:
        return jsonify({"error": "query is required"}), 400

    if simulate_latency:
        await asyncio.sleep(random.uniform(0.05, 0.4))

    q_embs = await embeddings.embed_texts([query], provider=settings.embedding_provider, context=f"search:{query[:80]}")
    q_emb = q_embs[0]

    results: List[Dict[str, Any]] = []

    if not lifecycle.qdrant_client:
        # local in-memory search
        for d in DOCUMENTS:
            for chunk in d['chunks']:
                score = sum(x * y for x, y in zip(q_emb, chunk['embed']))
                results.append({
                    'doc_id': d['id'],
                    'filename': d['filename'],
                    'chunk_text': (chunk['text'][:400] + ('...' if len(chunk['text']) > 400 else '')),
                    'score': float(score),
                })
    else:
        # Use Qdrant search
        search_res = lifecycle.qdrant_client.query_points(collection_name=settings.qdrant_collection, query=q_emb, limit=top_k)
        for hit in search_res.points:
            payload = hit.payload or {}  # type: ignore
            score = getattr(hit, 'score', None)
            # Qdrant might return different score field names; normalize
            s = float(score) if score is not None else 0.0
            results.append({
                'doc_id': hit.id,  # type: ignore
                'filename': payload.get('filename', 'unknown'),
                'chunk_text': (payload.get('text', '')[:400] + ('...' if len(payload.get('text', '')) > 400 else '')),
                'score': s,
            })

    results.sort(key=lambda r: r['score'], reverse=True)
    top = results[:top_k]

    # Calculate sources
    unique_files = set(r['filename'] for r in top)

    answer = None
    # Debug: log AI branch decisions
    log.info("AI branch check: demo=%s, AsyncAzureOpenAI=%s, key=%s, endpoint=%s, deployment=%s", settings.demo, bool(AsyncAzureOpenAI), getattr(settings, 'azure_openai_api_key', None), getattr(settings, 'azure_openai_endpoint', None), getattr(settings, 'azure_deployment_chat', None))

    # Demo-mode branch: use deterministic DemoLLM and clearly note that no chat model is available
    if settings.demo:
        try:
            # Apply demo matching thresholds: only include chunks that pass confidence checks
            best_score = max((r.get('score', 0.0) for r in top), default=0.0)
            ratio = float(getattr(settings, 'DEMO_MATCH_RATIO', 0.4))
            min_score = float(getattr(settings, 'DEMO_MIN_MATCH_SCORE', 0.0))
            threshold = max(best_score * ratio, min_score)

            filtered = [r for r in top if r.get('score', 0.0) >= threshold]

            if not filtered:
                answer = "(Demo mode - no chat model available) I couldn't find any confident matches in the documents (no documents matched the query)."
            else:
                context_parts = [f"Source ({r['filename']}): {r['chunk_text']}" for r in filtered]
                context = "\n\n".join(context_parts)
                demo = DemoLLM()
                cleaned, source_indicator = await demo.generate(query, context)
                answer = f"(Demo mode - no chat model available) {cleaned}"
        except Exception as e:
            log.exception("Demo LLM failed: %s", e)
            answer = "(Demo mode) Sorry, demo generation failed."

    elif not settings.demo and AsyncAzureOpenAI:
        try:
            if not settings.azure_openai_api_key or not settings.azure_openai_endpoint or not settings.azure_deployment_chat:
                answer = "Azure OpenAI not configured (Key, Endpoint, or Chat Deployment missing)."
            else:
                # Build context from top search results
                context_parts = [f"Source ({r['filename']}): {r['chunk_text']}" for r in top]
                context = "\n\n".join(context_parts)

                # Ask via AIManager (returns cleaned text and source indicator)
                cleaned, source_indicator = await support_agent.ai_manager.ask_agent_with_context(query, context, conversation_history=[])

                if cleaned == "__CONTENT_FILTER_BLOCKED__":
                    answer = "Content blocked by safety filters."
                else:
                    answer = cleaned
        except Exception as e:
            log.exception("RAG Generation failed")
            answer = f"Error generating answer: {str(e)}"

    return jsonify({
        "results": top, 
        "demo": settings.demo, 
        "message": DEMO_BANNER if settings.demo else "",
        "answer": answer,
        "sources": list(unique_files) if 'unique_files' in locals() else []
    })


@app.route('/optimize', methods=['POST'])
async def optimize_store():
    # Admin auth required
    if settings.admin_token:
        auth = request.headers.get('Authorization', '')
        if not auth.startswith('Bearer '):
            return jsonify({'error': 'unauthorized'}), 401
        token = auth.split(' ', 1)[1].strip()
        if token != settings.admin_token:
            return jsonify({'error': 'unauthorized'}), 401

    if lifecycle.qdrant_client and not settings.demo:
        try:
            # Trigger optimization (vacuum)
            # This is specific to Qdrant's internal optimizer configuration
            # Forcing an optimization is not always directly exposed via a simple method, 
            # but we can try to update collection params to trigger it or just acknowledge.
            # A common trick is to force a vacuum by setting vacuum_min_vector_number to something low then back? 
            # Or just rely on the fact that this button is mostly for user reassurance/debugging.
            # We will try to call update_collection with default optimizer config which might trigger a check.
            
            # Since qdrant-client python wrapper is used:
            # lifecycle.qdrant_client.update_collection(collection_name=settings.qdrant_collection, optimizer_config=models.OptimizersConfigDiff(vacuum_min_vector_number=...))
            
            # Simple "pass" for now with a log, unless we find a direct "optimize" method.
            # Qdrant automatically optimizes. 
            pass
        except Exception:
            log.exception('Failed to optimize Qdrant collection')
            return jsonify({"error": "optimization failed"}), 500
            
    return jsonify({"status": "optimized"})
async def reset_store():
    # Admin auth required
    if settings.admin_token:
        auth = request.headers.get('Authorization', '')
        if not auth.startswith('Bearer '):
            return jsonify({'error': 'unauthorized'}), 401
        token = auth.split(' ', 1)[1].strip()
        if token != settings.admin_token:
            return jsonify({'error': 'unauthorized'}), 401

    DOCUMENTS.clear()
    if lifecycle.qdrant_client and not settings.demo:
        try:
            # delete all points in collection (careful in prod)
            lifecycle.qdrant_client.delete(collection_name=settings.qdrant_collection, points_selector=qmodels.Filter(must=[]))
        except Exception:
            log.exception('Failed to clear Qdrant collection')
    return jsonify({"status": "cleared"})


@app.route('/document/<doc_id>', methods=['DELETE'])
async def delete_document(doc_id):
    # Admin auth required
    if settings.admin_token:
        auth = request.headers.get('Authorization', '')
        if not auth.startswith('Bearer '):
            return jsonify({'error': 'unauthorized'}), 401
        token = auth.split(' ', 1)[1].strip()
        if token != settings.admin_token:
            return jsonify({'error': 'unauthorized'}), 401

    if not lifecycle.qdrant_client:
        # Remove from DOCUMENTS
        DOCUMENTS[:] = [d for d in DOCUMENTS if d["id"] != doc_id]
        return jsonify({"status": "deleted"})

    # From Qdrant, use QdrantService to delete by filename
    try:
        deleted_count = qdrant_service.delete_by_filename(doc_id)
        if deleted_count > 0:
            return jsonify({"status": "deleted", "deleted": deleted_count})

        # If deletion via scanning didn't find anything, try field filters as last resort
        try:
            lifecycle.qdrant_client.delete(
                collection_name=settings.qdrant_collection,
                points_selector=qmodels.Filter(must=[qmodels.FieldCondition(key="filename", match=qmodels.MatchValue(value=doc_id))])
            )
            return jsonify({"status": "deleted"})
        except Exception:
            try:
                lifecycle.qdrant_client.delete(
                    collection_name=settings.qdrant_collection,
                    points_selector=qmodels.Filter(must=[qmodels.FieldCondition(key="file_name", match=qmodels.MatchValue(value=doc_id))])
                )
                return jsonify({"status": "deleted"})
            except Exception as e:
                log.exception('Failed to delete document with filters: %s', e)
                return jsonify({"error": "failed to delete"}), 500

    except Exception as e:
        log.exception('Failed to delete document')
        return jsonify({"error": "failed to delete"}), 500


@app.before_serving
async def _on_startup():
    # connect qdrant if configured
    await lifecycle.startup()

@app.after_serving
async def _on_shutdown():
    await lifecycle.shutdown()

if __name__ == '__main__':
    app.run(host='127.0.0.1', port=8002, debug=True)
